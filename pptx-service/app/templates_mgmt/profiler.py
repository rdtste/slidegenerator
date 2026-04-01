"""Deep template profiling — extracts comprehensive visual DNA from PowerPoint templates.

Consolidated from app.services.profile_service into the templates_mgmt module.

Extracts:
- Full color DNA (all scheme colors, chart color sequences)
- Typography DNA (fonts, sizes, weights across all placeholders)
- Per-layout visual details (picture areas, chart areas, text capacities)
- Design personality assessment
"""

from __future__ import annotations

import logging
import zipfile
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation

from app.templates_mgmt.models import (
    ColorDNA,
    TypographyDNA,
    LayoutDetail,
    PlaceholderDetail,
    ChartGuidelines,
    ImageGuidelines,
    TemplateProfile,
)
from app.templates_mgmt.service import get_template_path, load_presentation

logger = logging.getLogger(__name__)

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# Placeholder type constants
_PH_TITLE = 1
_PH_BODY = 2
_PH_SUBTITLE = 3
_PH_OBJECT = 7
_PH_TABLE = 10
_PH_CHART = 12
_PH_SLIDE_NUMBER = 13
_PH_DATE = 14
_PH_FOOTER = 15
_PH_DATE2 = 16
_PH_PICTURE = 18

_PH_SKIP = {_PH_SLIDE_NUMBER, _PH_DATE, _PH_FOOTER, _PH_DATE2}

_PH_TYPE_NAMES = {
    0: "TITLE_OR_CENTER", 1: "TITLE", 2: "BODY", 3: "SUBTITLE",
    7: "OBJECT", 10: "TABLE", 12: "CHART", 13: "SLIDE_NUMBER",
    14: "DATE", 15: "FOOTER", 16: "DATE", 18: "PICTURE",
}


def _emu_to_cm(emu: int) -> float:
    return round(emu / 914400 * 2.54, 2)


def extract_profile(template_id: str) -> TemplateProfile | None:
    """Extract a comprehensive visual profile from a template."""
    path = get_template_path(template_id)
    if not path:
        logger.warning(f"Template not found for profiling: {template_id}")
        return None

    profile = TemplateProfile(template_id=template_id)

    # 1. Extract colors + fonts from theme XML
    _extract_color_dna(path, profile)
    _extract_typography_dna(path, profile)

    # 2. Open the presentation for layout analysis
    try:
        prs = load_presentation(template_id)
    except Exception:
        logger.exception(f"Failed to load template for profiling: {template_id}")
        return profile

    profile.slide_width_cm = round(prs.slide_width / 914400 * 2.54, 1)
    profile.slide_height_cm = round(prs.slide_height / 914400 * 2.54, 1)

    # 3. Deep layout analysis
    _analyze_all_layouts(prs, profile)

    # 4. Build chart + image guidelines from extracted data
    _build_chart_guidelines(profile)
    _build_image_guidelines(profile)

    # 5. Determine supported layout types
    _determine_supported_types(profile)

    logger.info(
        f"Profile extracted for {template_id}: "
        f"{len(profile.layout_catalog)} layouts, "
        f"{len(profile.supported_layout_types)} supported types, "
        f"chart_colors={len(profile.color_dna.chart_colors)}"
    )
    return profile


def _extract_color_dna(path: Path, profile: TemplateProfile) -> None:
    """Parse the theme XML to extract all color scheme info including chart colors."""
    try:
        with open(path, "rb") as f:
            data = f.read()

        zf = zipfile.ZipFile(BytesIO(data))
        theme_files = sorted(
            n for n in zf.namelist()
            if "theme/theme" in n.lower() and n.endswith(".xml")
        )
        if not theme_files:
            return

        theme_xml = zf.read(theme_files[0])
        root = etree.fromstring(theme_xml)

        # Color scheme
        clr_scheme = root.find(".//a:clrScheme", _NS)
        if clr_scheme is not None:
            profile.template_name = clr_scheme.get("name", "")
            colors = parse_color_scheme(clr_scheme)

            dna = profile.color_dna
            dna.text = colors.get("dk1", dna.text)
            dna.background = colors.get("lt1", dna.background)
            dna.heading = colors.get("dk2", colors.get("dk1", dna.heading))
            dna.secondary = colors.get("lt2", dna.secondary)
            dna.accent1 = colors.get("accent1", dna.accent1)
            dna.accent2 = colors.get("accent2", dna.accent2)
            dna.accent3 = colors.get("accent3", dna.accent3)
            dna.accent4 = colors.get("accent4", dna.accent4)
            dna.accent5 = colors.get("accent5", dna.accent5)
            dna.accent6 = colors.get("accent6", dna.accent6)
            dna.primary = dna.accent1
            if "hlink" in colors:
                dna.hyperlink = colors["hlink"]

            # Chart color sequence = accent1 through accent6
            dna.chart_colors = [
                dna.accent1, dna.accent2, dna.accent3,
                dna.accent4, dna.accent5, dna.accent6,
            ]

        # Try to extract additional fill scheme colors for chart gradients
        fmt_scheme = root.find(".//a:fmtScheme", _NS)
        if fmt_scheme is not None:
            fill_styles = fmt_scheme.findall(".//a:fillStyleLst/a:solidFill/a:schemeClr", _NS)
            for fs in fill_styles:
                val = fs.get("val", "")
                logger.debug(f"Fill style scheme color: {val}")

    except Exception:
        logger.exception(f"Failed to extract color DNA from {path}")


def _extract_typography_dna(path: Path, profile: TemplateProfile) -> None:
    """Extract font information from theme XML."""
    try:
        with open(path, "rb") as f:
            data = f.read()

        zf = zipfile.ZipFile(BytesIO(data))
        theme_files = sorted(
            n for n in zf.namelist()
            if "theme/theme" in n.lower() and n.endswith(".xml")
        )
        if not theme_files:
            return

        theme_xml = zf.read(theme_files[0])
        root = etree.fromstring(theme_xml)

        typo = profile.typography_dna
        font_scheme = root.find(".//a:fontScheme", _NS)
        if font_scheme is not None:
            major = font_scheme.find(".//a:majorFont/a:latin", _NS)
            minor = font_scheme.find(".//a:minorFont/a:latin", _NS)
            if major is not None:
                typo.heading_font = major.get("typeface", typo.heading_font)
            if minor is not None:
                typo.body_font = minor.get("typeface", typo.body_font)

    except Exception:
        logger.exception(f"Failed to extract typography DNA from {path}")


def _analyze_all_layouts(prs: Presentation, profile: TemplateProfile) -> None:
    """Analyze every layout in the presentation for its visual capabilities."""
    heading_sizes: set[float] = set()
    body_sizes: set[float] = set()

    slide_w_cm = profile.slide_width_cm
    slide_h_cm = profile.slide_height_cm

    for idx, layout in enumerate(prs.slide_layouts):
        detail = LayoutDetail(index=idx, name=layout.name)

        content_phs = []
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in _PH_SKIP:
                continue

            type_name = _PH_TYPE_NAMES.get(ph_type, f"UNKNOWN_{ph_type}")
            detail.placeholder_types.append(type_name)

            width_cm = _emu_to_cm(ph.width) if ph.width else 0
            height_cm = _emu_to_cm(ph.height) if ph.height else 0
            left_cm = _emu_to_cm(ph.left) if ph.left else 0
            top_cm = _emu_to_cm(ph.top) if ph.top else 0

            # Classify spatial position based on coordinates
            position = _classify_position(left_cm, top_cm, width_cm, height_cm, slide_w_cm, slide_h_cm)

            font_sizes = _get_font_sizes_from_ph(ph)

            # Store detailed placeholder info
            ph_detail = PlaceholderDetail(
                type=type_name,
                index=ph.placeholder_format.idx,
                left_cm=left_cm,
                top_cm=top_cm,
                width_cm=width_cm,
                height_cm=height_cm,
                font_sizes_pt=font_sizes,
                position=position,
            )
            detail.placeholder_details.append(ph_detail)

            # Detect special placeholder types
            if ph_type == _PH_PICTURE:
                detail.has_picture = True
                detail.picture_width_cm = width_cm
                detail.picture_height_cm = height_cm
                if height_cm > 0:
                    ratio = width_cm / height_cm
                    detail.picture_aspect_ratio = _classify_aspect_ratio(ratio)

            elif ph_type == _PH_CHART:
                detail.has_chart = True

            elif ph_type == _PH_TABLE:
                detail.has_table = True

            # Track content area dimensions
            if ph_type in (_PH_OBJECT, _PH_BODY) and height_cm > 3:
                detail.content_width_cm = max(detail.content_width_cm, width_cm)
                detail.content_height_cm = max(detail.content_height_cm, height_cm)
                content_phs.append(ph)

            # Collect font sizes
            if ph_type == _PH_TITLE:
                heading_sizes.update(font_sizes)
            elif ph_type in (_PH_BODY, _PH_OBJECT):
                body_sizes.update(font_sizes)

        # Build a human-readable spatial description from placeholder details
        detail.spatial_description = _build_spatial_description(detail)

        # Compute text constraints for this layout
        _compute_layout_constraints(detail, layout)

        profile.layout_catalog.append(detail)

    # Update typography DNA
    if heading_sizes:
        profile.typography_dna.heading_sizes_pt = sorted(heading_sizes, reverse=True)
    if body_sizes:
        profile.typography_dna.body_sizes_pt = sorted(body_sizes, reverse=True)


def _classify_position(left_cm: float, top_cm: float, width_cm: float, height_cm: float,
                       slide_w: float, slide_h: float) -> str:
    """Classify a placeholder's spatial position relative to the slide."""
    center_x = left_cm + width_cm / 2
    coverage_w = width_cm / slide_w if slide_w > 0 else 0

    if coverage_w > 0.85:
        return "full-width"
    if center_x < slide_w * 0.35:
        return "left"
    if center_x > slide_w * 0.65:
        return "right"
    return "center"


def _build_spatial_description(detail: LayoutDetail) -> str:
    """Build a human-readable spatial description of how placeholders are arranged."""
    parts: list[str] = []
    for ph in detail.placeholder_details:
        size = f"{ph.width_cm:.0f}x{ph.height_cm:.0f}cm"
        parts.append(f"{ph.type}({ph.position}, {size})")
    return " | ".join(parts) if parts else ""


def _compute_layout_constraints(detail: LayoutDetail, layout) -> None:
    """Compute max_bullets, max_chars_per_bullet, title_max_chars from placeholder dimensions."""
    title_ph = None
    content_ph = None

    for ph in layout.placeholders:
        ph_type = ph.placeholder_format.type
        if ph_type in _PH_SKIP:
            continue

        if ph_type == _PH_TITLE and title_ph is None:
            title_ph = ph
        elif ph_type in (_PH_OBJECT, _PH_BODY) and content_ph is None:
            width_cm = _emu_to_cm(ph.width) if ph.width else 0
            height_cm = _emu_to_cm(ph.height) if ph.height else 0
            if height_cm > 3:
                content_ph = ph

    if title_ph:
        font_sizes = _get_font_sizes_from_ph(title_ph)
        font = font_sizes[0] if font_sizes else 18
        title_width = _emu_to_cm(title_ph.width) if title_ph.width else 20
        detail.title_max_chars = max(10, int((title_width / (font * 0.022)) * 0.8))

    if content_ph:
        font_sizes = _get_font_sizes_from_ph(content_ph)
        font = font_sizes[0] if font_sizes else 18
        width_cm = _emu_to_cm(content_ph.width) if content_ph.width else 20
        height_cm = _emu_to_cm(content_ph.height) if content_ph.height else 10

        raw_lines = height_cm / (font * 0.065)
        detail.max_bullets = max(1, int(raw_lines * 0.8))
        detail.max_chars_per_bullet = max(10, int((width_cm / (font * 0.022)) * 0.8))


def _get_font_sizes_from_ph(ph) -> list[float]:
    """Extract font sizes from a placeholder's XML definition."""
    sizes: list[float] = []
    try:
        xml = ph._element
        for rPr in xml.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"):
            sz = rPr.get("sz")
            if sz:
                sizes.append(int(sz) / 100)
        for defRPr in xml.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr"):
            sz = defRPr.get("sz")
            if sz:
                sizes.append(int(sz) / 100)
    except Exception:
        pass
    return sorted(set(sizes)) if sizes else []


def _classify_aspect_ratio(ratio: float) -> str:
    """Classify a width/height ratio into a standard aspect ratio."""
    if ratio > 1.6:
        return "16:9"
    elif ratio > 1.2:
        return "4:3"
    elif ratio > 0.9:
        return "1:1"
    elif ratio > 0.7:
        return "3:4"
    else:
        return "9:16"


def _build_chart_guidelines(profile: TemplateProfile) -> None:
    """Build chart generation guidelines from the extracted profile."""
    dna = profile.color_dna
    typo = profile.typography_dna

    chart_layouts = [
        ld.index for ld in profile.layout_catalog
        if ld.has_chart
    ]

    profile.chart_guidelines = ChartGuidelines(
        color_sequence=dna.chart_colors,
        font_family=typo.body_font,
        font_size_pt=typo.body_sizes_pt[0] if typo.body_sizes_pt else 10.0,
        background_color="transparent",
        grid_color=lighten_color(dna.text, 0.8),
        text_color=dna.text,
        style="modern_flat",
        available_chart_layouts=chart_layouts,
    )


def _build_image_guidelines(profile: TemplateProfile) -> None:
    """Build image generation guidelines from the extracted profile."""
    image_layouts = []
    primary_ratio = "16:9"

    for ld in profile.layout_catalog:
        if ld.has_picture:
            image_layouts.append(ld.index)
            if ld.picture_aspect_ratio:
                primary_ratio = ld.picture_aspect_ratio

    profile.image_guidelines = ImageGuidelines(
        available_image_layouts=image_layouts,
        primary_aspect_ratio=primary_ratio,
        style_keywords=["professional", "corporate", "modern", "clean"],
        accent_color=profile.color_dna.accent1,
    )


def _determine_supported_types(profile: TemplateProfile) -> None:
    """Determine which abstract layout types this template can support based on its layouts."""
    types: set[str] = set()

    for ld in profile.layout_catalog:
        ph_types = set(ld.placeholder_types)

        # Every template supports these basics at minimum
        if "TITLE" in ph_types or "BODY" in ph_types:
            types.add("title")
            types.add("section")
            types.add("closing")

        if "OBJECT" in ph_types:
            types.add("content")

        if ld.has_picture:
            types.add("image")

        if ld.has_chart:
            types.add("chart")

        if ld.has_table:
            types.add("table")

        # Detect two_column: two OBJECT placeholders
        object_count = ld.placeholder_types.count("OBJECT")
        if object_count >= 2:
            types.add("two_column")

    # Always include base types
    types.update({"title", "section", "content", "closing"})

    profile.supported_layout_types = sorted(types)


def parse_color_scheme(clr_scheme) -> dict[str, str]:
    """Parse a clrScheme element into a dict of {name: #hex}."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    colors: dict[str, str] = {}
    for child in clr_scheme:
        tag = child.tag.split("}")[-1]
        srgb = child.find("a:srgbClr", ns)
        sys_clr = child.find("a:sysClr", ns)
        if srgb is not None:
            colors[tag] = f"#{srgb.get('val', '000000')}"
        elif sys_clr is not None:
            colors[tag] = f"#{sys_clr.get('lastClr', sys_clr.get('val', '000000'))}"
    return colors


def lighten_color(hex_color: str, factor: float) -> str:
    """Lighten a hex color by mixing with white."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return "#E0E0E0"
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02x}{g:02x}{b:02x}"
