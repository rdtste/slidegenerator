"""PPTX generation service — maps structured slide data onto PowerPoint master layouts."""

from __future__ import annotations

import asyncio
import json
import logging
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from contextvars import ContextVar
from typing import Callable, Optional

from app.models.schemas import PresentationData, SlideContent
from app.services import template_service
from app.services.image_service import create_fallback_image, generate_image, generate_image_async
from app.services.chart_service import generate_chart, parse_chart_data
from app.services.image_fitting import fit_image_to_placeholder

logger = logging.getLogger(__name__)

ProgressCallback = Callable[[str, str, Optional[int]], None]
_progress_ctx: ContextVar[Optional[ProgressCallback]] = ContextVar(
    "_progress_ctx", default=None
)
_warnings_ctx: ContextVar[list[dict] | None] = ContextVar("_warnings_ctx", default=None)
_current_slide_number_ctx: ContextVar[int | None] = ContextVar("_current_slide_number_ctx", default=None)
_current_slide_title_ctx: ContextVar[str] = ContextVar("_current_slide_title_ctx", default="")
# Pre-generated image cache: description → Path (filled before slide loop)
_prefetched_images: ContextVar[dict[str, Path | None]] = ContextVar("_prefetched_images", default={})
_current_title_limit_ctx: ContextVar[int | None] = ContextVar("_current_title_limit_ctx", default=None)
_current_template_id_ctx: ContextVar[str] = ContextVar("_current_template_id_ctx", default="default")

_LAYOUT_LABELS: dict[str, str] = {
    "title": "Titelfolie",
    "section": "Kapitelfolie",
    "content": "Inhalt",
    "two_column": "Zwei Spalten",
    "image": "Bildfolie",
    "chart": "Diagramm",
    "closing": "Abschlussfolie",
}


def _report_progress(step: str, message: str, progress: int | None = None) -> None:
    cb = _progress_ctx.get(None)
    if cb is not None:
        cb(step, message, progress)


def _report_warning(message: str, code: str = "generation_warning") -> None:
    """Send a visible generation warning to SSE and optionally collect it for response payloads."""
    slide_number = _current_slide_number_ctx.get(None)
    slide_title = _current_slide_title_ctx.get("")
    payload = {
        "code": code,
        "message": message,
        "slide": slide_number,
        "title": slide_title,
    }
    warnings = _warnings_ctx.get(None)
    if warnings is not None:
        warnings.append(payload)
    logger.warning(
        "Generation warning slide=%s title=%s code=%s message=%s",
        slide_number,
        slide_title,
        code,
        message,
    )
    _report_progress("warning", message, None)

# Placeholder type constants from the OOXML spec
_PH_TITLE = 1       # TITLE
_PH_BODY = 2        # BODY
_PH_SUBTITLE = 3    # SUBTITLE (rare in custom templates)
_PH_OBJECT = 7      # OBJECT (generic content)
_PH_PICTURE = 18    # PICTURE

# Scored keyword patterns per layout type.
# Each entry: (substring_to_match, score, negative_substrings_that_disqualify)
_LAYOUT_SCORES: dict[str, list[tuple[str, int, list[str]]]] = {
    "title": [
        ("titelfolie", 80, ["bild", "farbig", "foto"]),
        ("titelfolie", 50, ["farbig"]),
        ("title slide", 70, []),
        ("titel", 30, ["nur", "farbig", "kapitelbeginn"]),
    ],
    "section": [
        ("kapitelbeginn", 90, ["inhalt"]),
        ("kapitelbeginn", 60, []),
        ("abschnitt", 70, []),
        ("section", 70, []),
        ("zwischentitel", 60, []),
    ],
    "content": [
        ("inhalt 1-spaltig", 100, []),
        ("inhalt 1", 90, []),
        ("content", 70, ["picture", "contact"]),
        ("inhalt", 50, ["spaltig", "bild", "kapitel"]),
        ("aufz\u00e4hlung", 60, []),
        ("bullet", 60, []),
        ("text", 30, ["nur"]),
    ],
    "two_column": [
        ("inhalt 2-spaltig", 100, []),
        ("2-spaltig", 90, []),
        ("two column", 80, []),
        ("vergleich", 70, []),
        ("comparison", 70, []),
    ],
    "image": [
        ("bild + inhalt", 90, []),
        ("bild + box", 80, []),
        ("bild + inhalt (1:2)", 85, []),
        ("bild", 60, ["titelfolie"]),
        ("image", 60, []),
        ("picture", 50, []),
        ("foto", 50, []),
    ],
    "chart": [
        ("diagramm", 90, []),
        ("chart", 90, []),
        ("grafik", 70, ["hintergrund"]),
        ("bild + inhalt", 40, []),
        ("bild", 30, ["titelfolie"]),
    ],
    "closing": [
        ("danke", 100, []),
        ("thank", 100, []),
        ("abschluss", 90, []),
        ("ende", 80, []),
        ("closing", 80, []),
        ("kontakt", 40, ["team", "bild"]),
        ("titelfolie", 20, ["bild", "farbig", "foto"]),
    ],
}

# Fallback layout indices for standard PowerPoint templates (no custom template)
_FALLBACK_LAYOUT: dict[str, int] = {
    "title": 0,
    "section": 2,
    "content": 1,
    "two_column": 3,
    "image": 5,
    "chart": 5,
    "closing": 0,
}


def generate_pptx(
    data: PresentationData,
    template_id: str = "default",
    progress_callback: ProgressCallback | None = None,
    warnings_collector: list[dict] | None = None,
    custom_color: str | None = None,
    custom_font: str | None = None,
) -> Path:
    """Generate a PPTX file from structured presentation data."""
    token = _progress_ctx.set(progress_callback)
    warnings_token = _warnings_ctx.set(warnings_collector if warnings_collector is not None else [])
    prefetch_token = _prefetched_images.set({})
    template_token = _current_template_id_ctx.set(template_id)
    try:
        _report_progress("template", "Template wird geladen...", 5)
        prs = template_service.load_presentation(template_id)

        analysis_map = _load_analysis_map(template_id)
        title_limits = _load_title_limits(template_id)
        _remove_all_slides(prs)
        _set_metadata(prs, data)
        _report_progress("template", "Template bereit", 10)

        # --- Pre-generate all images in parallel ---
        image_descs = _collect_image_descriptions(data)
        if image_descs:
            _report_progress("image", f"{len(image_descs)} Bilder werden parallel generiert…", 12)
            prefetched = _prefetch_images(image_descs)
            _prefetched_images.set(prefetched)
            ok_count = sum(1 for v in prefetched.values() if v is not None)
            _report_progress("image", f"{ok_count}/{len(image_descs)} Bilder bereit", 40)

        # Pre-render content leak sanitization — same quality gate as V2
        from app.validators.v1_content_leak_check import sanitize_presentation
        sanitize_presentation(data)

        total = len(data.slides)
        for i, slide_data in enumerate(data.slides):
            progress = 40 + int((i / total) * 50)
            label = _LAYOUT_LABELS.get(slide_data.layout, slide_data.layout)
            title_preview = f" — {slide_data.title[:40]}" if slide_data.title else ""
            _report_progress(
                "slide", f"Folie {i + 1}/{total}: {label}{title_preview}", progress,
            )
            _add_slide(prs, slide_data, analysis_map, title_limits, i + 1)
            _report_progress(
                "slide", f"Folie {i + 1}/{total} fertig", 40 + int(((i + 1) / total) * 50),
            )

        # Apply custom design (color/font) for default template
        if (not template_id or template_id == "default") and (custom_color or custom_font):
            _apply_custom_design(prs, custom_color, custom_font)

        _report_progress("saving", "Präsentation wird gespeichert...", 95)
        output_dir = Path(tempfile.mkdtemp(prefix="slidegen_"))
        safe_title = "".join(
            c if c.isalnum() or c in " _-" else "_" for c in data.title
        )[:50]
        output_path = output_dir / f"{safe_title}.pptx"
        prs.save(str(output_path))

        logger.info(f"Generated PPTX: {output_path} ({len(data.slides)} slides)")
        return output_path
    finally:
        _current_template_id_ctx.reset(template_token)
        _prefetched_images.reset(prefetch_token)
        _warnings_ctx.reset(warnings_token)
        _progress_ctx.reset(token)


def _apply_custom_design(prs: Presentation, custom_color: str | None, custom_font: str | None) -> None:
    """Apply custom accent color and font to all slides (default template only)."""
    color = None
    if custom_color:
        try:
            hex_val = custom_color.lstrip("#")
            color = RGBColor(int(hex_val[:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))
        except (ValueError, IndexError):
            logger.warning(f"Invalid custom color: {custom_color}")

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if custom_font:
                        run.font.name = custom_font
                    # Apply accent color to titles (placeholder types 0=title, 1=center_title, 13=title)
                    if color and hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx in (0, 1, 13):  # Title placeholders
                            run.font.color.rgb = color

    logger.info(f"Applied custom design: color={custom_color}, font={custom_font}")


def _collect_image_descriptions(data: PresentationData) -> list[str]:
    """Collect all image descriptions that will need generation.

    For image slides, enriches the description with slide title and bullet
    context so the generated image is relevant to the actual slide content.
    """
    descs: list[str] = []
    seen: set[str] = set()
    for slide_data in data.slides:
        desc = None
        if slide_data.layout == "title":
            desc = slide_data.image_description or slide_data.subtitle or slide_data.title
        elif slide_data.layout == "image":
            # Build a rich description from all slide content for relevance
            desc = _build_image_context(slide_data)
        if desc and desc not in seen:
            descs.append(desc)
            seen.add(desc)
    return descs


def _build_image_context(slide_data: SlideContent) -> str:
    """Build a contextual image description from slide title, bullets, and image_description.

    Combines the explicit image alt-text with the slide's actual content
    so the generated image matches what the audience reads on the slide.
    """
    parts: list[str] = []

    # Start with the explicit image description if available
    if slide_data.image_description:
        parts.append(slide_data.image_description)

    # Add title context
    if slide_data.title and slide_data.title not in (slide_data.image_description or ""):
        parts.append(f"Topic: {slide_data.title}")

    # Add bullet context (key messages on the slide)
    if slide_data.bullets:
        bullet_summary = "; ".join(b[:60] for b in slide_data.bullets[:4])
        parts.append(f"Context: {bullet_summary}")

    return ". ".join(parts) if parts else slide_data.body or ""


def _prefetch_images(descriptions: list[str]) -> dict[str, Path | None]:
    """Generate all images in parallel using asyncio."""
    style_kw = _load_image_style_keywords() or None
    async def _run():
        tasks = [generate_image_async(desc, style_keywords=style_kw) for desc in descriptions]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        cache: dict[str, Path | None] = {}
        for desc, result in zip(descriptions, results):
            if isinstance(result, Exception):
                logger.warning(f"[Prefetch] Image failed for '{desc[:40]}': {result}")
                cache[desc] = None
            else:
                cache[desc] = result
        return cache

    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # We're inside an async context (e.g. FastAPI) — use new loop in thread
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, _run())
                return future.result(timeout=120)
        return loop.run_until_complete(_run())
    except RuntimeError:
        return asyncio.run(_run())


def _get_prefetched_image(desc: str) -> Path | None:
    """Get a pre-generated image from the cache, or generate on demand as fallback."""
    cache = _prefetched_images.get({})
    if desc in cache:
        return cache[desc]
    # Fallback: generate synchronously (shouldn't happen normally)
    return generate_image(desc)


def _remove_all_slides(prs: Presentation) -> None:
    """Remove all existing slides from a presentation, keeping only layouts/masters."""
    slide_count = len(prs.slides)
    if slide_count == 0:
        return

    sldIdLst = prs.slides._sldIdLst
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    for _ in range(slide_count):
        sldId = sldIdLst[0]
        rId = sldId.get(f"{ns}id") or sldId.get("r:id")
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    logger.info(f"Removed {slide_count} existing template slides")


def _set_metadata(prs: Presentation, data: PresentationData) -> None:
    """Set presentation metadata."""
    prs.core_properties.title = data.title
    if data.author:
        prs.core_properties.author = data.author


def _add_slide(
    prs: Presentation,
    slide_data: SlideContent,
    analysis_map: dict[str, int] | None = None,
    title_limits: dict[int, int] | None = None,
    slide_number: int | None = None,
) -> None:
    """Add a single slide to the presentation based on its layout type."""
    layout_idx = _resolve_layout(prs, slide_data.layout, analysis_map)
    layout = prs.slide_layouts[layout_idx]

    # Closing slides with substantive content need adequate space.
    # Contact-style closing layouts often have tiny content areas (< 12cm wide).
    # Fall back to a content layout if the closing layout is too cramped.
    if slide_data.layout == "closing" and slide_data.bullets and len(slide_data.bullets) > 0:
        content_width_cm = _max_content_width_cm(layout)
        if content_width_cm < 12.0:
            content_idx = _resolve_layout(prs, "content", analysis_map)
            logger.info(
                f"Closing layout [{layout_idx}] '{layout.name}' content area too narrow "
                f"({content_width_cm:.1f}cm) for {len(slide_data.bullets)} bullets — "
                f"using content layout [{content_idx}] instead"
            )
            layout_idx = content_idx
            layout = prs.slide_layouts[layout_idx]

    slide = prs.slides.add_slide(layout)

    configured_title_limit = (title_limits or {}).get(layout_idx)
    estimated_title_limit = _estimate_title_capacity(layout)
    effective_title_limit = configured_title_limit if (configured_title_limit and configured_title_limit > 0) else estimated_title_limit

    token_slide_num = _current_slide_number_ctx.set(slide_number)
    token_slide_title = _current_slide_title_ctx.set(slide_data.title)
    token_title_limit = _current_title_limit_ctx.set(effective_title_limit)

    logger.debug(
        f"Added slide: layout={slide_data.layout} -> "
        f"idx={layout_idx} ({layout.name}), "
        f"title_limit={effective_title_limit}, "
        f"placeholders={[ph.placeholder_format.idx for ph in slide.placeholders]}"
    )

    try:
        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes

        handler = _LAYOUT_HANDLERS.get(slide_data.layout, _handle_content)
        handler(slide, slide_data)
    finally:
        _current_title_limit_ctx.reset(token_title_limit)
        _current_slide_title_ctx.reset(token_slide_title)
        _current_slide_number_ctx.reset(token_slide_num)


def _load_analysis_map(template_id: str) -> dict[str, int] | None:
    """Load AI-generated layout type → index mapping from analysis JSON."""
    template_path = template_service.get_template_path(template_id)
    if not template_path:
        return None

    analysis_path = template_path.parent / f"{template_id}.analysis.json"
    if not analysis_path.is_file():
        return None

    try:
        with open(analysis_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        mappings = data.get("layout_mappings", [])
        result: dict[str, int] = {}
        for m in mappings:
            mapped_type = m.get("mapped_type", "")
            idx = m.get("layout_index")
            if mapped_type and mapped_type != "unused" and isinstance(idx, int):
                # First mapping per type wins (AI returns best-match first)
                if mapped_type not in result:
                    result[mapped_type] = idx
        if result:
            logger.info(f"Loaded AI analysis map for '{template_id}': {result}")
        return result if result else None
    except Exception:
        logger.warning(f"Failed to load analysis for '{template_id}', using keyword fallback")
        return None


def _load_title_limits(template_id: str) -> dict[int, int]:
    """Load per-layout title character limits from profile/analysis metadata."""
    template_path = template_service.get_template_path(template_id)
    if not template_path:
        return {}

    candidates = [
        template_path.parent / f"{template_id}.analysis.json",
        template_path.parent / f"{template_id}.profile.json",
    ]

    for path in candidates:
        if not path.is_file():
            continue
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)

            mappings = data.get("layout_mappings") or data.get("layout_catalog") or []
            limits: dict[int, int] = {}
            for m in mappings:
                idx = m.get("layout_index")
                max_chars = m.get("title_max_chars")
                if isinstance(idx, int) and isinstance(max_chars, int) and max_chars > 0:
                    limits[idx] = max_chars
            if limits:
                logger.info(f"Loaded title limits from {path.name}: {limits}")
            return limits
        except Exception:
            logger.warning(f"Failed to parse title limits from {path.name}")

    return {}


def _resolve_layout(prs: Presentation, layout_type: str, analysis_map: dict[str, int] | None = None) -> int:
    """Resolve layout type to a template slide layout index.

    Uses AI analysis mapping when available, falls back to scored keyword matching.
    """
    # 1. Try AI analysis mapping
    if analysis_map and layout_type in analysis_map:
        idx = analysis_map[layout_type]
        if 0 <= idx < len(prs.slide_layouts) and _layout_supports_type(prs.slide_layouts[idx], layout_type):
            logger.info(
                f"Layout '{layout_type}' -> [{idx}] "
                f'"{prs.slide_layouts[idx].name}" (AI analysis)'
            )
            return idx
        logger.warning(
            f"Layout '{layout_type}' mapped to [{idx}] by analysis, "
            "but placeholder structure does not match. Falling back to scoring."
        )

    # 2. Fall back to scored keyword matching
    score_rules = _LAYOUT_SCORES.get(layout_type)
    if not score_rules:
        fallback = _FALLBACK_LAYOUT.get(layout_type, 1)
        return min(fallback, len(prs.slide_layouts) - 1)

    best_idx = -1
    best_score = -10_000

    for idx, layout in enumerate(prs.slide_layouts):
        name_lower = layout.name.lower()
        structure_bonus = _structure_score(layout, layout_type)
        layout_score = structure_bonus
        for pattern, score, negatives in score_rules:
            if pattern in name_lower:
                if any(neg in name_lower for neg in negatives):
                    continue
                layout_score += score
                break  # first matching pattern wins for this layout

        if layout_score > best_score:
            best_score = layout_score
            best_idx = idx

    if best_idx >= 0 and best_score > -100:
        logger.info(
            f"Layout '{layout_type}' -> [{best_idx}] "
            f'"{prs.slide_layouts[best_idx].name}" (score={best_score})'
        )
        return best_idx

    fallback = _FALLBACK_LAYOUT.get(layout_type, 1)
    result = min(fallback, len(prs.slide_layouts) - 1)
    logger.warning(f"Layout '{layout_type}' -> fallback [{result}]")
    return result


def _placeholder_type_counts(layout) -> dict[int, int]:
    counts: dict[int, int] = {}
    for ph in layout.placeholders:
        t = ph.placeholder_format.type
        counts[t] = counts.get(t, 0) + 1
    return counts


def _layout_supports_type(layout, layout_type: str) -> bool:
    counts = _placeholder_type_counts(layout)
    content_like = counts.get(_PH_OBJECT, 0) + counts.get(_PH_BODY, 0)

    if layout_type == "image":
        return counts.get(_PH_PICTURE, 0) > 0
    if layout_type == "chart":
        return counts.get(_PH_PICTURE, 0) > 0 or content_like > 0
    if layout_type == "two_column":
        return content_like >= 2
    if layout_type == "title":
        return counts.get(_PH_TITLE, 0) > 0 or content_like > 0
    if layout_type in {"content", "section", "closing"}:
        return content_like > 0 or counts.get(_PH_TITLE, 0) > 0
    return True


def _structure_score(layout, layout_type: str) -> int:
    counts = _placeholder_type_counts(layout)
    content_like = counts.get(_PH_OBJECT, 0) + counts.get(_PH_BODY, 0)
    has_picture = counts.get(_PH_PICTURE, 0) > 0

    if layout_type == "image":
        score = 160 if has_picture else -120
        if content_like > 0:
            score += 40
        return score
    if layout_type == "chart":
        score = 120 if has_picture else 0
        if content_like > 0:
            score += 60
        return score
    if layout_type == "two_column":
        return 90 if content_like >= 2 else -60
    if layout_type == "content":
        score = 70 if content_like > 0 else -40
        if has_picture:
            score -= 20
        return score
    if layout_type == "title":
        return 60 if (counts.get(_PH_TITLE, 0) > 0 or content_like > 0) else -40
    if layout_type in {"section", "closing"}:
        return 40 if (counts.get(_PH_TITLE, 0) > 0 or content_like > 0) else -30
    return 0


def _max_content_width_cm(layout) -> float:
    """Return the width in cm of the widest OBJECT/BODY placeholder in a layout."""
    max_w = 0.0
    for ph in layout.placeholders:
        if ph.placeholder_format.type in (_PH_OBJECT, _PH_BODY):
            w_cm = ph.width / 914400 * 2.54
            if w_cm > max_w:
                max_w = w_cm
    return max_w


def _estimate_title_capacity(layout) -> int:
    """Estimate max title chars from title/body placeholder width when profile limits are missing."""
    candidate = None
    for ph in layout.placeholders:
        if ph.placeholder_format.type == _PH_TITLE:
            candidate = ph
            break

    if candidate is None:
        for ph in layout.placeholders:
            if ph.placeholder_format.type in (_PH_BODY, _PH_OBJECT):
                candidate = ph
                break

    if candidate is None:
        return 60

    width_in = max(1.0, candidate.width / 914400)
    estimated = int(width_in * 8.5)
    return max(40, min(120, estimated))


# --- Layout handlers ---
# All handlers use _find_ph_by_type() to locate placeholders by their
# semantic type rather than by hardcoded index.

def _handle_title(slide, data: SlideContent) -> None:
    """Populate a title slide."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    body_phs = _find_all_ph_by_types(slide, [_PH_BODY, _PH_SUBTITLE, _PH_OBJECT])

    if title_ph:
        title_ph.text = _truncate_title(data.title)
        if data.subtitle and body_phs:
            body_phs[0].text = data.subtitle
    elif body_phs:
        # No TITLE placeholder (e.g. REWE template): use first BODY for title, second for subtitle
        body_phs[0].text = _truncate_title(data.title)
        if data.subtitle and len(body_phs) > 1:
            body_phs[1].text = data.subtitle

    # Fill picture placeholder if the title layout has one (e.g. "Titelfolie + Bild")
    picture_ph = _find_ph_by_type(slide, _PH_PICTURE)
    if picture_ph:
        desc = data.image_description or data.subtitle or data.title
        if desc:
            ph_width = picture_ph.width
            ph_height = picture_ph.height
            width_px = max(512, min(1536, int(ph_width / 914400 * 96)))
            height_px = max(512, min(1536, int(ph_height / 914400 * 96)))
            image_path = _get_prefetched_image(desc)
            if image_path is None:
                _report_warning(
                    "Bild konnte nicht per KI erzeugt werden - Platzhalterbild wird verwendet.",
                    code="image_generation_failed",
                )
                image_path = create_fallback_image(desc, width=width_px, height=height_px)
            if image_path:
                try:
                    fitted = fit_image_to_placeholder(image_path, ph_width, ph_height)
                    picture_ph.insert_picture(str(fitted))
                    logger.info(f"Inserted title image: {desc[:60]}...")
                except Exception:
                    logger.exception("Failed to insert title image")


def _handle_section(slide, data: SlideContent) -> None:
    """Populate a section header slide, optionally with content bullets."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    body_phs = _find_all_ph_by_types(slide, [_PH_BODY, _PH_SUBTITLE])
    content_ph = _find_content_placeholder(slide)

    if title_ph:
        title_ph.text = _truncate_title(data.title)
        if data.subtitle and body_phs:
            body_phs[0].text = data.subtitle
    elif body_phs:
        body_phs[0].text = _truncate_title(data.title)
        if data.subtitle and len(body_phs) > 1:
            body_phs[1].text = data.subtitle

    # Fill OBJECT placeholder with bullets if available (e.g. "Kapitelbeginn + Inhalt")
    if content_ph and data.bullets:
        _fill_bullet_list(content_ph, data.bullets)


def _handle_content(slide, data: SlideContent) -> None:
    """Populate a content slide with title and bullets."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    content_ph = _find_content_placeholder(slide)

    if title_ph:
        title_ph.text = _truncate_title(data.title)
    else:
        # Fallback: use a BODY placeholder for the title
        body_phs = _find_all_ph_by_types(slide, [_PH_BODY])
        if body_phs:
            body_phs[0].text = _truncate_title(data.title)

    if content_ph is None:
        return

    if data.bullets:
        _fill_bullet_list(content_ph, data.bullets)
    elif data.body:
        _set_text_with_bold(content_ph, data.body)


def _handle_two_column(slide, data: SlideContent) -> None:
    """Populate a two-column slide."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    if title_ph:
        title_ph.text = _truncate_title(data.title)

    # Find the two OBJECT/BODY placeholders for columns
    col_phs = _find_all_ph_by_types(slide, [_PH_OBJECT])
    if len(col_phs) < 2:
        # Fallback to BODY types (skip the first if it's being used for title-like content)
        col_phs = _find_all_ph_by_types(slide, [_PH_OBJECT, _PH_BODY])
        if title_ph is None and len(col_phs) > 2:
            col_phs = col_phs[1:]  # skip the one used for title

    left_ph = col_phs[0] if len(col_phs) > 0 else None
    right_ph = col_phs[1] if len(col_phs) > 1 else None

    if left_ph and data.left_column:
        _fill_bullets(left_ph, data.left_column)
    elif left_ph and data.bullets:
        half = len(data.bullets) // 2 or 1
        _fill_bullet_list(left_ph, data.bullets[:half])

    if right_ph and data.right_column:
        _fill_bullets(right_ph, data.right_column)
    elif right_ph and data.bullets:
        half = len(data.bullets) // 2 or 1
        _fill_bullet_list(right_ph, data.bullets[half:])


def _handle_image(slide, data: SlideContent) -> None:
    """Populate an image slide: generate image and insert into picture placeholder.

    Also fills the OBJECT/content placeholder with bullets when the layout
    offers one (e.g. "Bild + Inhalt").
    """
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    if title_ph:
        title_ph.text = _truncate_title(data.title)

    picture_ph = _find_ph_by_type(slide, _PH_PICTURE)
    # Use the same contextual description as _collect_image_descriptions for cache lookup
    desc = _build_image_context(data)

    image_inserted = False
    if picture_ph and desc:
        _report_progress("image", f"Bild wird eingefügt: {data.title or desc[:60]}")
        ph_width = picture_ph.width
        ph_height = picture_ph.height
        width_px = max(512, min(1536, int(ph_width / 914400 * 96)))
        height_px = max(512, min(1536, int(ph_height / 914400 * 96)))

        image_path = _get_prefetched_image(desc)
        if image_path is None:
            image_path = create_fallback_image(desc, width=width_px, height=height_px)
        if image_path:
            try:
                fitted = fit_image_to_placeholder(image_path, ph_width, ph_height)
                picture_ph.insert_picture(str(fitted))
                logger.info(f"Inserted generated image into picture placeholder: {desc[:60]}...")
                image_inserted = True
            except Exception:
                logger.exception("Failed to insert image into picture placeholder")
    elif desc:
        _report_warning(
            "Layout hat keinen Bild-Platzhalter - Bild konnte nicht eingebettet werden.",
            code="missing_picture_placeholder",
        )

    # Fill the content placeholder with bullets/body text (for "Bild + Inhalt" layouts)
    # Design-Prinzip "Klarheit vor Dekoration": Jedes Element muss einen Zweck haben.
    content_ph = _find_content_placeholder(slide)
    if content_ph:
        if data.bullets:
            _fill_bullet_list(content_ph, data.bullets)
        elif data.body:
            _set_text_with_bold(content_ph, data.body)
        else:
            # No bullets and no body — leave content empty rather than showing image description
            logger.warning("Image slide '%s' has no text content — content area left empty", data.title)


def _handle_closing(slide, data: SlideContent) -> None:
    """Populate a closing slide with title, subtitle, and optional bullets/body."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    content_ph = _find_content_placeholder(slide)
    body_phs = _find_all_ph_by_types(slide, [_PH_BODY, _PH_SUBTITLE])

    if title_ph:
        title_ph.text = _truncate_title(data.title)
        if data.subtitle and body_phs:
            body_phs[0].text = data.subtitle
    elif body_phs:
        body_phs[0].text = data.title
        if data.subtitle and len(body_phs) > 1:
            body_phs[1].text = data.subtitle

    # Fill content/OBJECT placeholder with bullets or body text
    if content_ph:
        if data.bullets:
            _fill_bullet_list(content_ph, data.bullets)
        elif data.body:
            _set_text_with_bold(content_ph, data.body)


def _handle_chart(slide, data: SlideContent) -> None:
    """Populate a chart slide: generate chart image and insert into picture placeholder."""
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    if title_ph:
        title_ph.text = _truncate_title(data.title)

    # Parse chart JSON from the slide data
    chart_data = None
    if data.chart_data:
        import json
        try:
            chart_data = json.loads(data.chart_data)
        except (json.JSONDecodeError, ValueError):
            logger.warning(f"Failed to parse chart JSON: {data.chart_data[:100]}")

    if not chart_data:
        # Fallback: try to extract from body text
        chart_data = parse_chart_data(data.body) if data.body else None

    if not chart_data:
        # No chart data — fall back to content handler
        logger.warning("No chart data found, falling back to content handler")
        _handle_content(slide, data)
        return

    # Load template profile for chart styling
    chart_style = _load_chart_style()

    # Find picture or chart placeholder for the chart image
    picture_ph = _find_ph_by_type(slide, _PH_PICTURE)
    if not picture_ph:
        # Fall back to any OBJECT placeholder
        picture_ph = _find_content_placeholder(slide)

    if picture_ph:
        _report_progress("chart", f"Diagramm wird erstellt: {chart_data.get('type', 'chart')}")
        ph_width = picture_ph.width
        ph_height = picture_ph.height
        width_px = max(600, min(1800, int(ph_width / 914400 * 96)))
        height_px = max(400, min(1200, int(ph_height / 914400 * 96)))

        chart_path = generate_chart(
            chart_data=chart_data,
            colors=chart_style.get("colors"),
            font_family=chart_style.get("font_family", "Calibri"),
            text_color=chart_style.get("text_color", "#333333"),
            grid_color=chart_style.get("grid_color", "#E0E0E0"),
            width_px=width_px,
            height_px=height_px,
        )

        if chart_path:
            try:
                fitted = fit_image_to_placeholder(chart_path, ph_width, ph_height)
                if hasattr(picture_ph, 'insert_picture'):
                    picture_ph.insert_picture(str(fitted))
                else:
                    # OBJECT placeholder: add picture as a shape
                    from pptx.util import Emu
                    slide.shapes.add_picture(
                        str(fitted),
                        picture_ph.left, picture_ph.top,
                        picture_ph.width, picture_ph.height,
                    )
                logger.info(f"Inserted chart image: {chart_data.get('type', 'unknown')} chart")
                return
            except Exception:
                logger.exception("Failed to insert chart image")

    # Last resort: leave content area empty — never show internal metadata as text
    logger.warning("Chart could not be rendered and no fallback available for slide '%s'", data.title)


def _load_template_profile() -> dict | None:
    """Load the current template's .profile.json for chart/image styling."""
    import json

    template_id = _current_template_id_ctx.get()
    if not template_id or template_id == "default":
        return None

    template_path = template_service.get_template_path(template_id)
    if not template_path:
        return None

    profile_path = template_path.parent / f"{template_id}.profile.json"
    if not profile_path.is_file():
        return None

    try:
        with open(profile_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        logger.warning(f"Failed to load profile for '{template_id}'")
        return None


def _load_chart_style() -> dict:
    """Load chart styling (colors, font, text_color, grid_color) from the template profile."""
    style: dict = {}
    profile = _load_template_profile()
    if not profile:
        return style

    color_dna = profile.get("color_dna", {})
    chart_colors = color_dna.get("chart_colors", [])
    if chart_colors:
        style["colors"] = chart_colors

    chart_guide = profile.get("chart_guidelines", {})
    if chart_guide.get("font_family"):
        style["font_family"] = chart_guide["font_family"]
    if chart_guide.get("text_color"):
        style["text_color"] = chart_guide["text_color"]
    if chart_guide.get("grid_color"):
        style["grid_color"] = chart_guide["grid_color"]

    return style


def _load_image_style_keywords() -> list[str]:
    """Load image style keywords from the template profile for Imagen prompt enrichment."""
    profile = _load_template_profile()
    if not profile:
        return []

    image_guide = profile.get("image_guidelines", {})
    keywords = image_guide.get("style_keywords", [])
    accent = image_guide.get("accent_color", "")

    # Add accent color as style hint
    if accent and accent != "#0969da":
        keywords = list(keywords) + [f"accent color {accent}"]

    return keywords


_LAYOUT_HANDLERS = {
    "title": _handle_title,
    "section": _handle_section,
    "content": _handle_content,
    "two_column": _handle_two_column,
    "image": _handle_image,
    "chart": _handle_chart,
    "closing": _handle_closing,
}


# --- Placeholder helpers ---

def _find_ph_by_type(slide, ph_type: int):
    """Find the first placeholder matching a given type constant."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            return ph
    return None


def _find_all_ph_by_types(slide, ph_types: list[int]) -> list:
    """Find all placeholders matching any of the given types, ordered by idx.

    Excludes footer/date/slide-number placeholders (types 13, 15, 16).
    """
    _SKIP_TYPES = {13, 15, 16}  # SLIDE_NUMBER, FOOTER, DATE
    result = []
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if t in _SKIP_TYPES:
            continue
        if t in ph_types:
            result.append(ph)
    result.sort(key=lambda p: p.placeholder_format.idx)
    return result


def _find_content_placeholder(slide):
    """Find the best placeholder for body/bullet content.

    Prefers OBJECT type, then BODY type. Skips the TITLE placeholder.
    """
    # First try OBJECT placeholders (most common for content in custom templates)
    for ph in slide.placeholders:
        if ph.placeholder_format.type == _PH_OBJECT:
            return ph

    # Then try BODY placeholders, skipping any that look like title/subtitle
    title_ph = _find_ph_by_type(slide, _PH_TITLE)
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if t == _PH_BODY and ph != title_ph:
            # Skip date/footer/slide-number
            if ph.placeholder_format.idx >= 10:
                return ph

    # Last resort: any BODY placeholder
    for ph in slide.placeholders:
        if ph.placeholder_format.type == _PH_BODY:
            return ph

    return None


import re as _re

# Regex to split text on **bold** markers
_BOLD_RE = _re.compile(r'\*\*(.+?)\*\*')


def _fill_bullets(placeholder, markdown_text: str) -> None:
    """Fill a placeholder with bullet lines parsed from Markdown.

    Supports indented sub-bullets (level 1) via 2+ leading spaces or tabs.
    """
    items: list[tuple[int, str]] = []
    for line in markdown_text.strip().split("\n"):
        if not line.strip():
            continue
        # Detect indentation level: 2+ spaces or tab = sub-bullet (level 1)
        stripped = line.lstrip()
        indent = len(line) - len(stripped)
        level = 1 if indent >= 2 else 0
        text = stripped.lstrip("- ").lstrip("* ").strip()
        if text:
            items.append((level, text))
    _fill_bullet_list_leveled(placeholder, items)


_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Default bullet character used when the template has buNone at level 0
_DEFAULT_BULLET_CHAR = "\u2022"  # •
_DEFAULT_BULLET_FONT = "Arial"


def _needs_bullet_char(tf) -> bool:
    """Check if the text frame's level-0 paragraphs lack a bullet character.

    Inspects the lstStyle for buNone or missing buChar at level 0/1.
    When the master defines buNone, bullets render as plain text — we fix that.
    """
    txBody = tf._txBody
    lstStyle = txBody.find(f"{_NS_A}lstStyle")
    if lstStyle is not None and len(lstStyle) > 0:
        lvl1 = lstStyle.find(f"{_NS_A}lvl1pPr")
        if lvl1 is not None:
            if lvl1.find(f"{_NS_A}buNone") is not None:
                return True
            if lvl1.find(f"{_NS_A}buChar") is not None:
                return False
    # Fall back: check if placeholder inherits from layout/master with buNone
    # If no bullet info at all, default to adding bullets for OBJECT placeholders
    return True


def _ensure_bullet_char(paragraph) -> None:
    """Add an explicit bullet character to a paragraph's XML properties."""
    pPr = paragraph._p.find(f"{_NS_A}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f"{_NS_A}pPr")
        paragraph._p.insert(0, pPr)

    # Remove buNone if present
    buNone = pPr.find(f"{_NS_A}buNone")
    if buNone is not None:
        pPr.remove(buNone)

    # Add buFont + buChar if not already there
    if pPr.find(f"{_NS_A}buChar") is None:
        buFont = etree.SubElement(pPr, f"{_NS_A}buFont")
        buFont.set("typeface", _DEFAULT_BULLET_FONT)
        buChar = etree.SubElement(pPr, f"{_NS_A}buChar")
        buChar.set("char", _DEFAULT_BULLET_CHAR)


def _truncate_title(text: str, max_chars: int | None = None) -> str:
    """Truncate a title to fit within the maximum character limit."""
    if not text:
        return text

    resolved_limit = max_chars if max_chars and max_chars > 0 else _current_title_limit_ctx.get(None)
    if not resolved_limit or resolved_limit <= 0:
        resolved_limit = 50

    if len(text) <= resolved_limit:
        return text
    # Try to cut at a word boundary
    truncated = text[:resolved_limit - 1]  # -1 to leave room for ellipsis
    last_space = truncated.rfind(" ")
    if last_space > resolved_limit * 0.6:
        truncated = truncated[:last_space]
    truncated = truncated.rstrip(" ,:;-") + "…"
    slide_title = _current_slide_title_ctx.get("")
    logger.warning(f"Title truncated: '{text[:60]}' → '{truncated}' (limit {resolved_limit})")
    return truncated


def _safe_clear_text_frame(tf):
    """Clear text frame content while preserving the template's lstStyle element.

    python-pptx's tf.clear() removes the lstStyle which contains font sizes,
    bullet symbols and indentation inherited from the slide layout / master.
    """
    from copy import deepcopy

    txBody = tf._txBody
    lstStyle = txBody.find(f"{_NS_A}lstStyle")
    lstStyle_copy = deepcopy(lstStyle) if lstStyle is not None and len(lstStyle) > 0 else None

    tf.clear()
    tf.word_wrap = True

    if lstStyle_copy is not None:
        new_lstStyle = txBody.find(f"{_NS_A}lstStyle")
        if new_lstStyle is not None:
            txBody.replace(new_lstStyle, lstStyle_copy)
        else:
            first_p = txBody.find(f"{_NS_A}p")
            if first_p is not None:
                txBody.insert(list(txBody).index(first_p), lstStyle_copy)


def _fill_bullet_list(placeholder, items: list[str]) -> None:
    """Fill a placeholder with a list of bullet strings (all level 0)."""
    _fill_bullet_list_leveled(placeholder, [(0, item) for item in items])


def _fill_bullet_list_leveled(placeholder, items: list[tuple[int, str]]) -> None:
    """Fill a placeholder with leveled bullet strings, rendering **bold** as actual bold.

    Each item is (level, text) where level 0 = main bullet, level 1 = sub-bullet.
    Preserves the template's lstStyle to maintain font sizes, bullet symbols and indentation.
    Adds explicit bullet characters when the template doesn't provide them.
    Adds paragraph spacing for visual breathing room between bullets.
    """
    tf = placeholder.text_frame
    _safe_clear_text_frame(tf)

    needs_explicit_bullet = _needs_bullet_char(tf)

    # Calculate spacing based on number of items.
    n = max(len(items), 1)
    if n <= 3:
        space_after_pt = Pt(14)
    elif n <= 5:
        space_after_pt = Pt(10)
    elif n <= 7:
        space_after_pt = Pt(6)
    else:
        space_after_pt = Pt(4)

    for i, (level, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = min(level, 1)
        # Main bullets get full spacing; sub-bullets get tighter spacing
        p.space_after = Pt(2) if level > 0 else space_after_pt
        # Sub-bullets before a main bullet: reduce the main bullet's space_before
        if needs_explicit_bullet:
            _ensure_bullet_char(p)
        _set_paragraph_with_bold(p, text)
        # Slightly smaller font for sub-bullets to create visual hierarchy
        if level > 0 and p.runs:
            for run in p.runs:
                if run.font.size is None:
                    run.font.size = Pt(14)
                else:
                    run.font.size = Pt(max(10, int(run.font.size / 12700) - 2))

    # Remove trailing space on last paragraph
    if tf.paragraphs:
        tf.paragraphs[-1].space_after = Pt(0)

    # Enable shrink-to-fit so text auto-scales when exceeding the placeholder
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_paragraph_with_bold(paragraph, text: str) -> None:
    """Parse markdown bold (**text**) and create proper bold/normal runs."""
    # Split text on **bold** markers
    parts = _BOLD_RE.split(text)
    # parts alternates: [normal, bold, normal, bold, ...]
    # Even indices are normal text, odd indices are bold text
    if len(parts) == 1 and '**' not in text:
        # No bold formatting — simple text
        paragraph.text = text
        return

    # Clear any existing runs, then add formatted runs
    paragraph.clear()
    for idx, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        if idx % 2 == 1:
            # Odd index = was inside **bold** markers
            run.font.bold = True


def _set_text_with_bold(placeholder, text: str) -> None:
    """Set placeholder text, rendering any **bold** markdown as actual bold.
    
    Preserves the template's lstStyle to maintain font sizes.
    """
    tf = placeholder.text_frame
    _safe_clear_text_frame(tf)
    _set_paragraph_with_bold(tf.paragraphs[0], text)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
