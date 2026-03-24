"""Template service — manages PowerPoint master templates."""

from __future__ import annotations

import logging
import shutil
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation

from app.config import settings
from app.models.schemas import TemplateInfo

logger = logging.getLogger(__name__)

# python-pptx cannot open POTX files directly because they have a different
# content type. We patch the content type inside the ZIP to make it work.
_POTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def _template_dirs() -> list[Path]:
    """Return template search directories in priority order.

    Priority:
    1. Configured PPTX service templates directory
    2. Local backend/templates directory (dev fallback)
    """
    dirs: list[Path] = [settings.templates_dir]

    # In local development, templates are often stored in backend/templates.
    backend_templates = Path(__file__).resolve().parents[3] / "backend" / "templates"
    if backend_templates.exists() and backend_templates not in dirs:
        dirs.append(backend_templates)

    return dirs


def _load_potx_as_presentation(path: Path) -> Presentation:
    """Load a .potx file by patching the content type in-memory."""
    with open(path, "rb") as f:
        data = f.read()

    buf = BytesIO(data)
    out = BytesIO()

    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(out, "w") as zout:
        for item in zin.infolist():
            content = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                content = content.replace(
                    _POTX_CONTENT_TYPE.encode(),
                    _PPTX_CONTENT_TYPE.encode(),
                )
            zout.writestr(item, content)

    out.seek(0)
    return Presentation(out)


def list_templates() -> list[TemplateInfo]:
    """List all available PowerPoint templates."""
    settings.templates_dir.mkdir(parents=True, exist_ok=True)

    templates: list[TemplateInfo] = []
    seen_ids: set[str] = set()
    for directory in _template_dirs():
        if not directory.exists():
            continue
        for pattern in ("*.pptx", "*.potx"):
            for path in sorted(directory.glob(pattern)):
                if path.stem in seen_ids:
                    continue
                seen_ids.add(path.stem)
            info = _inspect_template(path)
            if info:
                templates.append(info)

    if not templates:
        templates.append(_get_default_template_info())

    return templates


def get_template_path(template_id: str) -> Path | None:
    """Get the file path for a template by its ID."""
    for directory in _template_dirs():
        for ext in (".pptx", ".potx"):
            path = directory / f"{template_id}{ext}"
            if path.is_file():
                return path
    return None


def load_presentation(template_id: str) -> Presentation:
    """Load a PowerPoint presentation from template. Falls back to blank."""
    path = get_template_path(template_id)
    if path:
        logger.info(f"Loading template: {path.name}")
        if path.suffix.lower() == ".potx":
            return _load_potx_as_presentation(path)
        return Presentation(str(path))

    logger.info("No template found, using blank presentation")
    return Presentation()


def get_layout_names(prs: Presentation) -> list[str]:
    """Get the names of all slide layouts in a presentation."""
    return [layout.name for layout in prs.slide_layouts]


def find_layout(prs: Presentation, keywords: list[str]) -> int:
    """Find a slide layout index by matching keywords against layout names.

    Returns the index of the first matching layout, or 0 as fallback.
    """
    layout_names = [layout.name.lower() for layout in prs.slide_layouts]

    for keyword in keywords:
        for idx, name in enumerate(layout_names):
            if keyword.lower() in name:
                return idx

    return 0


def _inspect_template(path: Path) -> TemplateInfo | None:
    """Inspect a .pptx/.potx template and extract metadata."""
    try:
        if path.suffix.lower() == ".potx":
            prs = _load_potx_as_presentation(path)
        else:
            prs = Presentation(str(path))
        layouts = [layout.name for layout in prs.slide_layouts]
        core = prs.core_properties
        name = core.title or path.stem.replace("_", " ").replace("-", " ").title()
        description = core.subject or f"Template mit {len(layouts)} Layouts"

        return TemplateInfo(
            id=path.stem,
            name=name,
            description=description,
            layouts=layouts,
        )
    except Exception:
        logger.exception(f"Failed to inspect template: {path}")
        return None


def _get_default_template_info() -> TemplateInfo:
    """Return info for the built-in blank template."""
    prs = Presentation()
    layouts = [layout.name for layout in prs.slide_layouts]
    return TemplateInfo(
        id="default",
        name="Standard (Blank)",
        description="Standard-PowerPoint-Template ohne Branding",
        layouts=layouts,
    )
