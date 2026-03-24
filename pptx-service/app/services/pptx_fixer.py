"""PPTX Fixer — applies programmatic corrections to generated PPTX files."""

from __future__ import annotations

import logging
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

from app.services.gemini_vision_qa import SlideIssue

logger = logging.getLogger(__name__)

_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_PH_PICTURE = 18


class FixResult:
    """Result of applying fixes to a PPTX file."""

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.fixes_applied: list[dict] = []
        self.fixes_failed: list[dict] = []
        self.changed_slides: list[int] = []

    @property
    def any_changes(self) -> bool:
        return len(self.fixes_applied) > 0


def apply_fixes(pptx_path: str, issues: list[SlideIssue]) -> FixResult:
    """Apply programmatic fixes to a PPTX file based on QA issues.

    Creates a backup, applies fixes, saves in-place.

    Args:
        pptx_path: Path to the PPTX file to fix.
        issues: List of issues with fix_action instructions.

    Returns:
        FixResult with details of applied fixes.
    """
    result = FixResult(pptx_path)
    fixable = [i for i in issues if i.fix_action != "none"]

    if not fixable:
        return result

    # Backup original
    backup_path = pptx_path + ".backup"
    shutil.copy2(pptx_path, backup_path)

    try:
        prs = Presentation(pptx_path)
        slides = list(prs.slides)

        for issue in fixable:
            slide_idx = issue.slide_number - 1
            if slide_idx < 0 or slide_idx >= len(slides):
                result.fixes_failed.append({
                    "slide": issue.slide_number,
                    "action": issue.fix_action,
                    "reason": "Folie nicht gefunden",
                })
                continue

            slide = slides[slide_idx]
            applied = _apply_single_fix(slide, issue, prs)

            if applied:
                result.fixes_applied.append({
                    "slide": issue.slide_number,
                    "action": issue.fix_action,
                    "element": issue.element,
                    "description": issue.description,
                })
                if issue.slide_number not in result.changed_slides:
                    result.changed_slides.append(issue.slide_number)
            else:
                result.fixes_failed.append({
                    "slide": issue.slide_number,
                    "action": issue.fix_action,
                    "reason": "Fix konnte nicht angewendet werden",
                })

        if result.any_changes:
            prs.save(pptx_path)
            logger.info(
                f"[PPTX Fixer] Saved {len(result.fixes_applied)} fixes "
                f"to {len(result.changed_slides)} slides"
            )

        # Clean up backup if everything went well
        Path(backup_path).unlink(missing_ok=True)

    except Exception as e:
        logger.error(f"[PPTX Fixer] Error applying fixes: {e}")
        # Restore from backup
        if Path(backup_path).exists():
            shutil.copy2(backup_path, pptx_path)
            Path(backup_path).unlink(missing_ok=True)
        result.fixes_failed.append({
            "slide": 0,
            "action": "all",
            "reason": f"Fehler: {str(e)[:100]}",
        })

    return result


def _apply_single_fix(slide, issue: SlideIssue, prs: Presentation) -> bool:
    """Apply a single fix to a slide. Returns True if fix was applied."""
    action = issue.fix_action

    try:
        if action == "resize_image":
            return _fix_resize_image(slide, issue)
        elif action == "crop_image":
            return _fix_crop_image(slide, issue)
        elif action == "truncate_text":
            return _fix_truncate_text(slide, issue)
        elif action == "remove_placeholder":
            return _fix_remove_placeholder(slide, issue)
        elif action == "reposition":
            return _fix_reposition(slide, issue)
        elif action == "change_font_color":
            return _fix_font_color(slide, issue)
        elif action == "fill_content":
            return _fix_fill_content(slide, issue)
        elif action == "adjust_spacing":
            return _fix_spacing(slide, issue)
        else:
            logger.debug(f"[PPTX Fixer] Unknown action: {action}")
            return False
    except Exception as e:
        logger.warning(
            f"[PPTX Fixer] Fix failed on slide {issue.slide_number} "
            f"({action}): {e}"
        )
        return False


def _fix_resize_image(slide, issue: SlideIssue) -> bool:
    """Resize an image that's too large to fit within the slide bounds."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slide_width = slide.slide_layout.slide_master.slide_width
    slide_height = slide.slide_layout.slide_master.slide_height

    for shape in slide.shapes:
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER):
            if not hasattr(shape, "image"):
                continue

            changed = False

            # Check if image exceeds slide bounds
            right_edge = shape.left + shape.width
            bottom_edge = shape.top + shape.height

            if right_edge > slide_width:
                # Scale down proportionally
                overflow_ratio = slide_width / right_edge
                shape.width = int(shape.width * overflow_ratio * 0.95)
                shape.height = int(shape.height * overflow_ratio * 0.95)
                changed = True

            if bottom_edge > slide_height:
                overflow_ratio = slide_height / bottom_edge
                shape.width = int(shape.width * overflow_ratio * 0.95)
                shape.height = int(shape.height * overflow_ratio * 0.95)
                changed = True

            if changed:
                logger.info(
                    f"[PPTX Fixer] Resized image on slide {issue.slide_number}"
                )
                return True

    return False


def _fix_crop_image(slide, issue: SlideIssue) -> bool:
    """Apply cropping to make image fit its placeholder."""
    # Cropping in python-pptx uses the crop properties on picture shapes
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # If the shape's aspect ratio doesn't match placeholder,
            # crop to fit (center crop)
            if hasattr(shape, "crop_left"):
                # Already a picture — we can adjust crop
                # For now, just ensure it doesn't overflow the slide
                return _fix_resize_image(slide, issue)

    return False


def _fix_truncate_text(slide, issue: SlideIssue) -> bool:
    """Truncate text that overflows its placeholder."""
    element_lower = issue.element.lower() if issue.element else ""

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        is_title = "titel" in element_lower or "title" in element_lower

        if is_title:
            # Only fix title-like shapes
            ph_type = getattr(shape.placeholder_format, "type", None) if hasattr(shape, "placeholder_format") else None
            if ph_type not in (1, None):  # 1 = TITLE
                continue

            for para in tf.paragraphs:
                text = para.text
                if len(text) > 50:
                    # Truncate at word boundary
                    truncated = text[:50]
                    last_space = truncated.rfind(" ")
                    if last_space > 30:
                        truncated = truncated[:last_space]
                    truncated = truncated.rstrip(" ,:;-")

                    # Clear and rewrite
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = truncated
                    else:
                        para.text = truncated
                    logger.info(
                        f"[PPTX Fixer] Truncated title on slide {issue.slide_number}"
                    )
                    return True
        else:
            # Truncate bullet text
            for para in tf.paragraphs:
                text = para.text
                if len(text) > 100:
                    truncated = text[:90]
                    last_space = truncated.rfind(" ")
                    if last_space > 50:
                        truncated = truncated[:last_space]
                    truncated = truncated.rstrip(" ,:;-")

                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = truncated
                    else:
                        para.text = truncated
                    logger.info(
                        f"[PPTX Fixer] Truncated text on slide {issue.slide_number}"
                    )
                    return True

    return False


def _fix_remove_placeholder(slide, issue: SlideIssue) -> bool:
    """Remove or clear placeholder text like 'Titel hinzufuegen'."""
    placeholder_patterns = [
        "titel hinzufuegen", "titel hinzufügen",
        "text hier", "text hinzufuegen", "text hinzufügen",
        "untertitel hinzufuegen", "untertitel hinzufügen",
        "xxxx", "lorem ipsum",
        "click to add", "add title", "add text",
    ]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        full_text = tf.text.lower().strip()

        if any(p in full_text for p in placeholder_patterns):
            # Clear the text frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            logger.info(
                f"[PPTX Fixer] Cleared placeholder text on slide {issue.slide_number}: "
                f"'{full_text[:40]}'"
            )
            return True

    return False


def _fix_reposition(slide, issue: SlideIssue) -> bool:
    """Reposition overlapping elements."""
    # Basic overlap fix: ensure no shape extends beyond slide bounds
    slide_width = slide.slide_layout.slide_master.slide_width
    slide_height = slide.slide_layout.slide_master.slide_height
    min_margin = Emu(457200)  # 0.5 inch

    changed = False
    for shape in slide.shapes:
        # Push shapes back into slide bounds
        if shape.left < 0:
            shape.left = min_margin
            changed = True
        if shape.top < 0:
            shape.top = min_margin
            changed = True

        right = shape.left + shape.width
        if right > slide_width:
            shape.left = max(min_margin, slide_width - shape.width - min_margin)
            changed = True

        bottom = shape.top + shape.height
        if bottom > slide_height:
            shape.top = max(min_margin, slide_height - shape.height - min_margin)
            changed = True

    if changed:
        logger.info(f"[PPTX Fixer] Repositioned elements on slide {issue.slide_number}")
    return changed


def _fix_font_color(slide, issue: SlideIssue) -> bool:
    """Fix low contrast by changing font color."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    r, g, b = run.font.color.rgb
                    # If text is very light (near white), make it dark
                    if r > 200 and g > 200 and b > 200:
                        run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x37)
                        logger.info(
                            f"[PPTX Fixer] Darkened text on slide {issue.slide_number}"
                        )
                        return True
                    # If text is very dark on presumably dark background
                    if r < 30 and g < 30 and b < 30:
                        # Leave as-is — dark text is usually fine
                        pass
    return False


def _fix_fill_content(slide, issue: SlideIssue) -> bool:
    """Fill empty content areas — mostly clears leftover placeholder text."""
    return _fix_remove_placeholder(slide, issue)


def _fix_spacing(slide, issue: SlideIssue) -> bool:
    """Adjust spacing between elements."""
    # Reduce font size slightly if text is cramped
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        para_count = len(tf.paragraphs)

        if para_count > 5:
            # Too many paragraphs — reduce spacing
            for para in tf.paragraphs:
                para.space_after = Pt(2)
                para.space_before = Pt(1)
            logger.info(
                f"[PPTX Fixer] Reduced spacing on slide {issue.slide_number}"
            )
            return True

    return False
