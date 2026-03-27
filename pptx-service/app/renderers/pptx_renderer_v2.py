"""V2 PPTX Renderer — Professional presentation rendering engine.

Design principles:
  1. Typography hierarchy — strict scale from Hero (44pt) down to Caption (11pt)
  2. Auto-fitting — text that would overflow is automatically reduced in size
  3. Whitespace first — generous padding, spacing, and breathing room
  4. Visual polish — accent bars, rounded cards, clean shapes, proper alignment
  5. Scannable output — every slide graspable in under 8 seconds
"""

from __future__ import annotations

import logging
import math
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

from app.schemas.models import RenderInstruction, RenderElement, SlideType

logger = logging.getLogger(__name__)

ProgressCallback = Callable[[str, str, int | None], None]


# ═══════════════════════════════════════════════════════════════
# Typography System
# ═══════════════════════════════════════════════════════════════

@dataclass(frozen=True)
class TypeStyle:
    """Immutable typography definition for one semantic role."""
    size_pt: int
    bold: bool = False
    color: str = "#333333"
    line_spacing: float = 1.15
    space_before_pt: int = 0
    space_after_pt: int = 4
    min_size_pt: int = 10


# Strict hierarchy — each level visually distinct from the next
_TYPO: dict[str, TypeStyle] = {
    # Slide-level headings
    "hero":          TypeStyle(44, bold=True, color="#FFFFFF", line_spacing=1.1,
                               space_after_pt=10, min_size_pt=30),
    "headline":      TypeStyle(28, bold=True, color="#1a1a2e", line_spacing=1.1,
                               space_after_pt=6, min_size_pt=20),
    "subheadline":   TypeStyle(18, color="#6b7280", line_spacing=1.2,
                               space_after_pt=4, min_size_pt=14),
    # Body content
    "text":          TypeStyle(16, color="#374151", line_spacing=1.45,
                               space_after_pt=4, min_size_pt=12),
    "bullet":        TypeStyle(17, color="#374151", line_spacing=1.5,
                               space_before_pt=10, space_after_pt=6, min_size_pt=13),
    # Card system
    "card_icon":     TypeStyle(36, color="#374151", line_spacing=1.0, min_size_pt=28),
    "card_title":    TypeStyle(17, bold=True, color="#1a1a2e", line_spacing=1.15,
                               space_after_pt=6, min_size_pt=14),
    "card_body":     TypeStyle(13, color="#4b5563", line_spacing=1.4, min_size_pt=11),
    # KPI system
    "kpi_value":     TypeStyle(48, bold=True, color="#1a1a2e", line_spacing=1.0,
                               min_size_pt=32),
    "kpi_label":     TypeStyle(13, color="#6b7280", line_spacing=1.2, min_size_pt=10),
    "kpi_delta":     TypeStyle(16, bold=True, color="#22c55e", line_spacing=1.0,
                               min_size_pt=12),
    # Quote / statement
    "statement":     TypeStyle(32, bold=True, color="#1a1a2e", line_spacing=1.35,
                               space_after_pt=8, min_size_pt=22),
    "attribution":   TypeStyle(15, color="#9ca3af", line_spacing=1.2, min_size_pt=12),
    # Timeline
    "timeline_date": TypeStyle(12, bold=True, color="#2563EB", line_spacing=1.1,
                               min_size_pt=10),
    "timeline_title": TypeStyle(14, bold=True, color="#1a1a2e", line_spacing=1.15,
                                min_size_pt=11),
    "timeline_desc": TypeStyle(11, color="#6b7280", line_spacing=1.25, min_size_pt=9),
    # Process
    "step_number":   TypeStyle(20, bold=True, color="#FFFFFF", line_spacing=1.0,
                               min_size_pt=16),
    "step_title":    TypeStyle(15, bold=True, color="#1a1a2e", line_spacing=1.15,
                               min_size_pt=12),
    "step_desc":     TypeStyle(12, color="#4b5563", line_spacing=1.3, min_size_pt=10),
    # Agenda / closing
    "agenda_item":   TypeStyle(20, color="#374151", line_spacing=1.8, min_size_pt=16),
    "closing_title": TypeStyle(34, bold=True, color="#1a1a2e", line_spacing=1.1,
                               space_after_pt=10, min_size_pt=26),
    "closing_body":  TypeStyle(17, color="#374151", line_spacing=1.7,
                               space_before_pt=10, space_after_pt=6, min_size_pt=14),
    "contact":       TypeStyle(13, color="#9ca3af", line_spacing=1.3, min_size_pt=11),
    # Utility
    "caption":       TypeStyle(11, color="#9ca3af", line_spacing=1.3, min_size_pt=9),
    "label":         TypeStyle(11, bold=True, color="#6b7280", line_spacing=1.2,
                               min_size_pt=9),
}


def _typo(role: str) -> TypeStyle:
    """Get typography style for a role, with body fallback."""
    return _TYPO.get(role, _TYPO["text"])


# Role detection from element key/type
_KEY_TO_ROLE: dict[str, str] = {
    "headline": "headline", "subheadline": "subheadline",
    "statement": "statement", "attribution": "attribution",
    "bullet_area": "bullet", "body_area": "text",
    "agenda_list": "agenda_item",
    "takeaways": "closing_body", "takeaway_area": "text",
    "contact": "contact",
}


def _infer_role(element_type: str, style=None) -> str:
    """Infer typography role from element_type string."""
    if element_type in _TYPO:
        return element_type
    role = _KEY_TO_ROLE.get(element_type)
    if role:
        return role
    # Check for known prefixes
    if element_type.startswith("card_title"):
        return "card_title"
    if element_type.startswith("card_body"):
        return "card_body"
    if element_type.startswith("card_icon"):
        return "card_icon"
    if element_type.startswith("kpi_value"):
        return "kpi_value"
    if element_type.startswith("kpi_label"):
        return "kpi_label"
    if element_type.startswith("kpi_delta"):
        return "kpi_delta"
    if element_type.startswith("date_"):
        return "timeline_date"
    if element_type.startswith("entry_title"):
        return "timeline_title"
    if element_type.startswith("entry_desc"):
        return "timeline_desc"
    if element_type.startswith("step_num"):
        return "step_number"
    if element_type.startswith("step_title"):
        return "step_title"
    if element_type.startswith("step_desc"):
        return "step_desc"
    return "text"


# ═══════════════════════════════════════════════════════════════
# Text Measurement & Auto-Fitting
# ═══════════════════════════════════════════════════════════════

def _estimate_lines(text: str, font_size_pt: int, width_cm: float) -> int:
    """Estimate rendered line count for text in a box of given width."""
    if not text:
        return 0
    # Heuristic: at size N pt, about 72/N * 0.45 chars fit per cm
    # (calibrated for Calibri/Arial proportional fonts)
    chars_per_cm = max(1.0, 72.0 / (font_size_pt * 0.55))
    chars_per_line = max(1, int(width_cm * chars_per_cm))
    lines = 0
    for para in text.split("\n"):
        lines += max(1, math.ceil(len(para) / chars_per_line)) if para else 1
    return lines


def _estimate_height_cm(text: str, font_size_pt: int, width_cm: float,
                        line_spacing: float = 1.15) -> float:
    """Estimate rendered text height in cm."""
    n_lines = _estimate_lines(text, font_size_pt, width_cm)
    line_h_cm = font_size_pt * line_spacing / 72.0 * 2.54
    return n_lines * line_h_cm


def _auto_fit(text: str, base_size: int, width_cm: float, height_cm: float,
              line_spacing: float, min_size: int = 10) -> int:
    """Reduce font size step-by-step until text fits available height."""
    size = base_size
    while size > min_size:
        h = _estimate_height_cm(text, size, width_cm, line_spacing)
        if h <= height_cm * 0.92:  # 8% safety margin
            return size
        size -= 1
    return min_size


def _auto_fit_bullets(items: list[str], base_size: int, width_cm: float,
                      height_cm: float, line_spacing: float,
                      space_between_pt: int = 10,
                      min_size: int = 13) -> int:
    """Auto-fit font size for a bullet list."""
    size = base_size
    while size > min_size:
        total_h = 0.0
        space_cm = space_between_pt / 72.0 * 2.54
        for item in items:
            total_h += _estimate_height_cm(str(item), size, width_cm, line_spacing)
            total_h += space_cm
        if total_h <= height_cm * 0.92:
            return size
        size -= 1
    return min_size


# ═══════════════════════════════════════════════════════════════
# Color Utilities
# ═══════════════════════════════════════════════════════════════

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return RGBColor(0x33, 0x33, 0x33)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighten(hex_color: str, factor: float) -> str:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return "#e5e7eb"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, int(r + (255 - r) * factor))
    g = min(255, int(g + (255 - g) * factor))
    b = min(255, int(b + (255 - b) * factor))
    return f"#{r:02x}{g:02x}{b:02x}"


_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_VANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


# ═══════════════════════════════════════════════════════════════
# Template Layout Mapping
# ═══════════════════════════════════════════════════════════════

_REWE_LAYOUT_MAP: dict[SlideType, int] = {
    SlideType.TITLE_HERO:       1,
    SlideType.SECTION_DIVIDER:  3,
    SlideType.KEY_STATEMENT:    13,
    SlideType.BULLETS_FOCUSED:  5,
    SlideType.THREE_CARDS:      7,
    SlideType.KPI_DASHBOARD:    6,
    SlideType.IMAGE_TEXT_SPLIT: 8,
    SlideType.COMPARISON:       6,
    SlideType.TIMELINE:         5,
    SlideType.PROCESS_FLOW:     5,
    SlideType.CHART_INSIGHT:    5,
    SlideType.IMAGE_FULLBLEED:  11,
    SlideType.AGENDA:           2,
    SlideType.CLOSING:          17,
}


# ═══════════════════════════════════════════════════════════════
# Paragraph Builder
# ═══════════════════════════════════════════════════════════════

class _ParagraphBuilder:
    """Fluent builder for creating well-formatted paragraphs."""

    def __init__(self, p, font_family: str):
        self._p = p
        self._font = font_family

    def align(self, alignment: str) -> "_ParagraphBuilder":
        self._p.alignment = _ALIGN.get(alignment, PP_ALIGN.LEFT)
        return self

    def spacing(self, before_pt: int = 0, after_pt: int = 4,
                line_spacing_pt: int | None = None) -> "_ParagraphBuilder":
        if before_pt:
            self._p.space_before = Pt(before_pt)
        if after_pt:
            self._p.space_after = Pt(after_pt)
        if line_spacing_pt and line_spacing_pt > 0:
            self._p.line_spacing = Pt(line_spacing_pt)
        return self

    def add_run(self, text: str, size_pt: int, color: str = "#333333",
                bold: bool = False, italic: bool = False) -> "_ParagraphBuilder":
        run = self._p.add_run()
        run.text = text
        run.font.name = self._font
        run.font.size = Pt(size_pt)
        run.font.color.rgb = _hex_to_rgb(color)
        run.font.bold = bold
        run.font.italic = italic
        return self


# ═══════════════════════════════════════════════════════════════
# Main Renderer
# ═══════════════════════════════════════════════════════════════

class PptxRendererV2:
    """Professional-grade PPTX renderer with auto-fitting and typography system."""

    def __init__(self, accent_color: str = "#2563EB", font_family: str = "Calibri",
                 template_id: str | None = None):
        self.accent = accent_color
        self.accent_light = _lighten(accent_color, 0.85)
        self.accent_dark = self._darken(accent_color, 0.15)
        self.font = font_family
        self.template_id = template_id
        self._image_generator: Callable | None = None
        self._chart_generator: Callable | None = None

    def set_image_generator(self, fn: Callable) -> None:
        self._image_generator = fn

    def set_chart_generator(self, fn: Callable) -> None:
        self._chart_generator = fn

    # ── Main entry point ──────────────────────────────────────

    def render(
        self,
        instructions: list[RenderInstruction],
        output_dir: str | None = None,
        progress_callback: ProgressCallback | None = None,
    ) -> Path:
        prs = self._load_presentation()
        use_tpl = self.template_id is not None

        if use_tpl:
            self._remove_all_slides(prs)

        total = len(instructions)
        for i, instr in enumerate(instructions):
            if progress_callback:
                pct = int(10 + (i / total) * 80)
                progress_callback("rendering", f"Folie {i+1}/{total} wird gerendert...", pct)

            slide = self._add_slide(prs, instr, use_tpl)

            if not use_tpl:
                self._set_background(slide, instr.background_color)

            for element in instr.elements:
                try:
                    self._render_element(slide, element, instr)
                except Exception as exc:
                    logger.warning(f"Slide {i+1} element '{element.element_type}': {exc}")

        out_dir = Path(output_dir or tempfile.mkdtemp())
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / "presentation_v2.pptx"
        prs.save(str(out_path))

        if progress_callback:
            progress_callback("saved", "Praesentation gespeichert", 95)

        return out_path

    # ── Presentation management ───────────────────────────────

    def _load_presentation(self) -> Presentation:
        if self.template_id:
            try:
                from app.services.template_service import load_presentation
                prs = load_presentation(self.template_id)
                logger.info(f"Loaded template: {self.template_id} "
                            f"({len(prs.slide_layouts)} layouts)")
                return prs
            except Exception as exc:
                logger.warning(f"Template '{self.template_id}' failed, using blank: {exc}")
        prs = Presentation()
        prs.slide_width = Cm(25.4)
        prs.slide_height = Cm(19.05)
        return prs

    def _remove_all_slides(self, prs: Presentation) -> None:
        sldIdLst = prs.part._element.sldIdLst
        if sldIdLst is None:
            return
        ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        for sldId in list(sldIdLst):
            rId = sldId.get(f'{ns}id')
            if rId:
                prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

    def _add_slide(self, prs, instr, use_tpl):
        if use_tpl:
            idx = _REWE_LAYOUT_MAP.get(instr.slide_type, 16)
            if idx < len(prs.slide_layouts):
                return prs.slides.add_slide(prs.slide_layouts[idx])
            fallback = min(16, len(prs.slide_layouts) - 1)
            return prs.slides.add_slide(prs.slide_layouts[fallback])
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _set_background(self, slide, color: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(color)

    # ── Element routing ───────────────────────────────────────

    def _render_element(self, slide, el: RenderElement,
                        instr: RenderInstruction) -> None:
        t = el.element_type
        if t == "shape":
            self._render_shape(slide, el, instr)
        elif t == "image":
            self._render_image(slide, el)
        elif t == "chart":
            self._render_chart(slide, el)
        elif t == "bullets":
            self._render_bullets(slide, el)
        else:
            self._render_text(slide, el, instr)

    # ── Text rendering (auto-fitted, typography-aware) ────────

    def _render_text(self, slide, el: RenderElement,
                     instr: RenderInstruction) -> None:
        text = str(el.content) if el.content else ""
        if not text.strip():
            return

        pos = el.position
        style = el.style
        role = _infer_role(el.element_type, style)
        typo = _typo(role)

        # Resolve effective values: style overrides > typography defaults
        base_size = style.font_size_pt or typo.size_pt
        color = self._resolve_color(style.font_color or typo.color)
        bold = style.bold or typo.bold
        spacing = style.line_spacing if style.line_spacing != 1.0 else typo.line_spacing
        alignment = style.alignment or "left"
        v_align = style.vertical_alignment or "top"

        # Auto-fit: reduce font if text would overflow
        fitted = _auto_fit(text, base_size, pos.width_cm, pos.height_cm,
                           spacing, typo.min_size_pt)

        txBox = slide.shapes.add_textbox(
            Cm(pos.left_cm), Cm(pos.top_cm),
            Cm(pos.width_cm), Cm(pos.height_cm),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = _VANCHOR.get(v_align, MSO_ANCHOR.TOP)

        line_sp_pt = int(fitted * spacing) if spacing > 1.0 else None

        # Render with markdown bold support
        if "**" in text:
            self._write_markdown_paragraph(
                tf.paragraphs[0], text, fitted, color, bold,
                alignment, line_sp_pt, typo,
            )
        else:
            pb = _ParagraphBuilder(tf.paragraphs[0], self.font)
            pb.align(alignment)
            pb.spacing(before_pt=typo.space_before_pt,
                       after_pt=typo.space_after_pt,
                       line_spacing_pt=line_sp_pt)
            pb.add_run(text, fitted, color, bold)

    def _write_markdown_paragraph(self, p, text: str, size: int, color: str,
                                  default_bold: bool, alignment: str,
                                  line_sp_pt: int | None,
                                  typo: TypeStyle) -> None:
        """Render text with **bold** markdown spans."""
        pb = _ParagraphBuilder(p, self.font)
        pb.align(alignment)
        pb.spacing(before_pt=typo.space_before_pt,
                   after_pt=typo.space_after_pt,
                   line_spacing_pt=line_sp_pt)

        parts = text.split("**")
        for i, part in enumerate(parts):
            if not part:
                continue
            is_bold = (i % 2 == 1)
            pb.add_run(part, size, color, bold=is_bold)

    # ── Bullet rendering (auto-fitted, proper spacing) ────────

    def _render_bullets(self, slide, el: RenderElement) -> None:
        items = el.content if isinstance(el.content, list) else []
        if not items:
            return

        pos = el.position
        style = el.style
        typo = _typo("bullet")

        base_size = style.font_size_pt or typo.size_pt
        color = self._resolve_color(style.font_color or typo.color)
        spacing = style.line_spacing if style.line_spacing != 1.0 else typo.line_spacing

        # Auto-fit all bullets to available height
        fitted = _auto_fit_bullets(
            [str(item) for item in items],
            base_size, pos.width_cm, pos.height_cm,
            spacing, typo.space_before_pt, typo.min_size_pt,
        )

        txBox = slide.shapes.add_textbox(
            Cm(pos.left_cm), Cm(pos.top_cm),
            Cm(pos.width_cm), Cm(pos.height_cm),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        line_sp_pt = int(fitted * spacing) if spacing > 1.0 else None
        alignment = style.alignment or "left"

        for idx, item_text in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            pb = _ParagraphBuilder(p, self.font)
            pb.align(alignment)
            pb.spacing(
                before_pt=typo.space_before_pt if idx > 0 else 0,
                after_pt=typo.space_after_pt,
                line_spacing_pt=line_sp_pt,
            )

            item_str = str(item_text)
            if "**" in item_str:
                # Bold-prefix pattern: **Label** rest of text
                parts = item_str.split("**")
                for i, part in enumerate(parts):
                    if not part:
                        continue
                    pb.add_run(part, fitted, color, bold=(i % 2 == 1))
            else:
                pb.add_run(f"\u2022  {item_str}", fitted, color)

    # ── Shape rendering ───────────────────────────────────────

    def _render_shape(self, slide, el: RenderElement,
                      instr: RenderInstruction) -> None:
        pos = el.position
        content = el.content if isinstance(el.content, dict) else {}
        fill_color = self._resolve_color(content.get("fill", "#e5e7eb"))
        corner_r = content.get("corner_radius_cm", 0)

        if corner_r > 0:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Cm(pos.left_cm), Cm(pos.top_cm),
                Cm(pos.width_cm), Cm(pos.height_cm),
            )
            max_adj = min(corner_r / min(pos.width_cm, pos.height_cm), 0.5)
            shape.adjustments[0] = max_adj
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Cm(pos.left_cm), Cm(pos.top_cm),
                Cm(pos.width_cm), Cm(pos.height_cm),
            )

        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        shape.line.fill.background()

    # ── Image rendering ───────────────────────────────────────

    def _render_image(self, slide, el: RenderElement) -> None:
        pos = el.position
        content = el.content if isinstance(el.content, dict) else {}
        description = content.get("description", "")

        if not description:
            return

        image_path = None
        if self._image_generator:
            try:
                image_path = self._image_generator(description)
            except Exception as exc:
                logger.warning(f"Image generation failed: {exc}")

        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(
                str(image_path),
                Cm(pos.left_cm), Cm(pos.top_cm),
                Cm(pos.width_cm), Cm(pos.height_cm),
            )
        else:
            self._render_image_placeholder(slide, pos, description)

    def _render_image_placeholder(self, slide, pos, description: str) -> None:
        """Render a styled placeholder when image generation is unavailable."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(pos.left_cm), Cm(pos.top_cm),
            Cm(pos.width_cm), Cm(pos.height_cm),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1f, 0x24, 0x2b)
        shape.line.fill.background()
        shape.adjustments[0] = 0.02

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        pb = _ParagraphBuilder(tf.paragraphs[0], self.font)
        pb.align("center")
        pb.add_run(f"\U0001F5BC  {description[:70]}", 12, "#9ca3af")

    # ── Chart rendering ───────────────────────────────────────

    def _render_chart(self, slide, el: RenderElement) -> None:
        pos = el.position
        content = el.content if isinstance(el.content, dict) else {}

        if not content or not self._chart_generator:
            self._render_chart_placeholder(slide, pos)
            return

        try:
            chart_path = self._chart_generator(content, self.accent)
            if chart_path and Path(chart_path).exists():
                slide.shapes.add_picture(
                    str(chart_path),
                    Cm(pos.left_cm), Cm(pos.top_cm),
                    Cm(pos.width_cm), Cm(pos.height_cm),
                )
        except Exception as exc:
            logger.warning(f"Chart generation failed: {exc}")

    def _render_chart_placeholder(self, slide, pos) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(pos.left_cm), Cm(pos.top_cm),
            Cm(pos.width_cm), Cm(pos.height_cm),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xf3, 0xf4, 0xf6)
        shape.line.fill.background()
        shape.adjustments[0] = 0.02

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        pb = _ParagraphBuilder(tf.paragraphs[0], self.font)
        pb.align("center")
        pb.add_run("\U0001F4CA  Chart", 14, "#9ca3af")

    # ── Color resolution ──────────────────────────────────────

    def _resolve_color(self, color: str) -> str:
        if color == "accent":
            return self.accent
        if color == "accent_light":
            return self.accent_light
        if color == "accent_dark":
            return self.accent_dark
        return color

    @staticmethod
    def _darken(hex_color: str, factor: float) -> str:
        h = hex_color.lstrip("#")
        if len(h) != 6:
            return "#1a1a2e"
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = int(r * (1 - factor))
        g = int(g * (1 - factor))
        b = int(b * (1 - factor))
        return f"#{r:02x}{g:02x}{b:02x}"
